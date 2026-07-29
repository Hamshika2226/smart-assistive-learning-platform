from flask import Flask, render_template, request, redirect, url_for, session, jsonify, send_from_directory, send_file
import os, json, uuid, time, re
os.environ["PATH"] += os.pathsep + r"C:\Users\hamsh\Downloads\ffmpeg-8.1-essentials_build\ffmpeg-8.1-essentials_build\bin"
import edge_tts
import asyncio
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.secret_key = "eduLearn_secure_key_2026"
UPLOAD_FOLDER = os.path.join(os.getcwd(), 'uploads')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
DB_FILE = 'db.json'

def load_db():
    if not os.path.exists(DB_FILE):
        save_db({'users': {}, 'content': []})
    with open(DB_FILE) as f:
        return json.load(f)

def save_db(data):
    with open(DB_FILE, 'w') as f:
        json.dump(data, f, indent=2)

BRAILLE = {
    'a':'⠁','b':'⠃','c':'⠉','d':'⠙','e':'⠑','f':'⠋','g':'⠛','h':'⠓','i':'⠊','j':'⠚',
    'k':'⠅','l':'⠇','m':'⠍','n':'⠝','o':'⠕','p':'⠏','q':'⠟','r':'⠗','s':'⠎','t':'⠞',
    'u':'⠥','v':'⠧','w':'⠺','x':'⠭','y':'⠽','z':'⠵',' ':'⠀',',':'⠂','.':'⠲','?':'⠦',
    '!':'⠖',':':'⠒',';':'⠆','-':'⠤','1':'⠁','2':'⠃','3':'⠉','4':'⠙','5':'⠑',
    '6':'⠋','7':'⠛','8':'⠓','9':'⠊','0':'⠚'
}

def to_braille(text):
    out = []
    for ch in text[:300]:
        if ch.isupper():
            out.append('⠠' + BRAILLE.get(ch.lower(), ch))
        else:
            out.append(BRAILLE.get(ch, ch))
    return ''.join(out)

def extract_text_from_pdf(filepath):
    try:
        import PyPDF2
        text = []
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text.append(t.strip())
        return '\n'.join(text).strip()
    except Exception as e:
        return f'[PDF extraction failed: {e}]'

def extract_text_from_ppt(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text.append(shape.text.strip())
        return '\n'.join(text).strip()
    except Exception as e:
        return f'[PPT extraction failed: {e}]'

def extract_text_from_audio_video(filepath):

    try:
        import whisper
        from moviepy import VideoFileClip
        import os
        audio_path = filepath + ".mp3"
        video = VideoFileClip(filepath)
        video.audio.write_audiofile(
            audio_path,
            logger=None
        )
        model = whisper.load_model("small")
        result = model.transcribe(
            audio_path,
            task="translate",
            fp16=False
        )
        segments = result.get("segments", [])
        lines = []
        for seg in segments:
            start = int(seg["start"])
            minutes = start // 60
            seconds = start % 60
            timestamp = f"{minutes:02d}:{seconds:02d}"
            text = seg["text"].strip()
            if text:
                lines.append(
                    f"[{timestamp}] {text}"
                )
        final_transcript = "\n".join(lines)
        if os.path.exists(audio_path):
            os.remove(audio_path)
        return final_transcript.strip()
    except Exception as e:
        import traceback
        traceback.print_exc()
        print("TRANSCRIPTION ERROR:", e)
        return ""

def generate_smart_summary(text, title=""):
    """Generate a meaningful educational summary — not raw transcript extraction."""
    if not text:
        return ''
    clean = re.sub(r'\[\d{2}:\d{2}\]\s*', '', text)
    clean = re.sub(r'\s+', ' ', clean).strip()
    sentences = re.split(r'(?<=[.!?])\s+', clean)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 25]

    if not sentences:
        return clean[:300]
    noise_words = ['presented by', 'guided by', 'references', 'date:', 'dr.', 'professor',
                   'module', 'chapter', 'unit', 'slide', 'page', 'thank you', 'questions?']
    filtered = [s for s in sentences if not any(w in s.lower() for w in noise_words)]
    if not filtered:
        filtered = sentences
    summary_parts = []
    if len(filtered) >= 3:
        summary_parts.append(filtered[0])
        mid = filtered[len(filtered) // 2]
        if mid not in summary_parts:
            summary_parts.append(mid)
        last = filtered[-1]
        if last not in summary_parts:
            summary_parts.append(last)
    else:
        summary_parts = filtered[:3]

    joined = ' '.join(summary_parts)
    if len(joined) < 80 and len(filtered) > 3:
        joined = ' '.join(filtered[:4])

    return joined.strip()

def summarize_as_bullets(text):
    """Return HTML bullet-point summary for API endpoint."""
    if not text:
        return '<ul><li>No content to summarize.</li></ul>'

    clean = re.sub(r'\[\d{2}:\d{2}\]\s*', '', text)
    clean = re.sub(r'\s+', ' ', clean).strip()

    sentences = re.split(r'(?<=[.!?])\s+', clean)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 30]

    noise_words = ['presented by', 'guided by', 'references', 'date:', 'dr.', 'professor',
                   'module', 'chapter', 'slide', 'page', 'thank you']
    filtered = [s for s in sentences if not any(w in s.lower() for w in noise_words)]
    if not filtered:
        filtered = sentences

    if len(filtered) >= 3:
        points = [filtered[0], filtered[len(filtered) // 2], filtered[-1]]
    else:
        points = filtered[:3]
    summary = '<ul>'
    for p in points:
        summary += f'<li>{p.strip().capitalize()}.</li>'
    summary += '</ul>'
    return summary

#route___________________
@app.route('/')
def index():
    if 'uid' in session:
        return redirect(url_for('teacher') if session.get('role') == 'teacher' else url_for('student'))
    return render_template('index.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if 'uid' in session:
        return redirect(url_for('teacher') if session.get('role') == 'teacher' else url_for('student'))
    if request.method == 'POST':
        db = load_db()
        email = request.form['email']
        pw = request.form['password']
        user = next((u for u in db['users'].values() if u['email'] == email and u['password'] == pw), None)
        if user:
            session.permanent = True
            session.update({
                'uid': user['id'],
                'name': user['name'],
                'role': user['role'],
                'disability': user.get('disability', 'none')
            })
            return redirect(url_for('teacher') if user['role'] == 'teacher' else url_for('student'))
        return render_template('login.html', error='Invalid credentials')
    return render_template('login.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    if 'uid' in session:
        return redirect(url_for('teacher') if session.get('role') == 'teacher' else url_for('student'))
    if request.method == 'POST':
        db = load_db()
        uid = str(uuid.uuid4())[:6]
        role = request.form['role']
        user = {
            'id': uid,
            'name': request.form['name'],
            'email': request.form['email'],
            'password': request.form['password'],
            'role': role,
            'disability': request.form.get('disability', 'none') if role == 'student' else 'none'
        }
        db['users'][uid] = user
        save_db(db)
        session.permanent = True
        session.update({
            'uid': uid,
            'name': user['name'],
            'role': role,
            'disability': user.get('disability', 'none')
        })
        return redirect(url_for('teacher') if role == 'teacher' else url_for('student'))
    return render_template('register.html')

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('index'))

# ── Teacher ──────────────────────────────────────────────────────────────────

@app.route('/teacher')
def teacher():
    if session.get('role') != 'teacher':
        return redirect('/login')
    db = load_db()
    return render_template('teacher.html', content=db['content'], name=session['name'])

@app.route('/teacher/upload', methods=['POST'])
def upload():
    if session.get('role') != 'teacher':
        return redirect('/login')

    db = load_db()
    title = request.form.get('title', 'Untitled')
    ctype = request.form.get('type', 'notes')
    disability_category = request.form.get('disability_category', 'all')
    text = request.form.get('text', '').strip()
    youtube_link = request.form.get('youtube_link', '').strip()

    file = request.files.get('file')
    sign_file = request.files.get('sign_video')

    filename = ''
    sign_filename = ''
    extracted = ''

    if file and file.filename:
        filename = secure_filename(f"{uuid.uuid4().hex[:6]}_{file.filename}")
        fpath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(fpath)
        fname_lower = filename.lower()

        if fname_lower.endswith('.pdf'):
            extracted = extract_text_from_pdf(fpath)
        elif fname_lower.endswith(('.ppt', '.pptx')):
            extracted = extract_text_from_ppt(fpath)
        elif fname_lower.endswith(('.mp3', '.wav', '.mp4', '.avi', '.mov')):
            extracted = extract_text_from_audio_video(fpath)
        elif fname_lower.endswith('.txt'):
            with open(fpath, 'r', errors='ignore') as f:
                extracted = f.read().strip()

    if sign_file and sign_file.filename:
        sign_filename = secure_filename(f"{uuid.uuid4().hex[:6]}_{sign_file.filename}")
        sign_file.save(os.path.join(UPLOAD_FOLDER, sign_filename))

    # Normalise YouTube link to embed format
    if 'youtu.be/' in youtube_link:
        youtube_link = youtube_link.replace('youtu.be/', 'www.youtube.com/embed/')
    elif 'watch?v=' in youtube_link:
        youtube_link = youtube_link.replace('watch?v=', 'embed/')

    final_text = text
    if extracted and not extracted.startswith('['):
        final_text = (text + '\n\n' + extracted).strip() if text else extracted
    elif not text and extracted:
        final_text = extracted

   # Smart educational summary
    summary = generate_smart_summary(
        extracted if extracted else final_text,
        title
    )

    status = 'completed' if final_text else 'failed'

    # AUTO INCREMENT CONTENT ID
    next_id = 1

    if db['content']:

        try:

            ids = []

            for c in db['content']:

                try:
                    ids.append(int(c.get('id', 0)))
                except:
                    pass

            if ids:
                next_id = max(ids) + 1

        except:
            next_id = len(db['content']) + 1

    item = {

        'id': str(next_id),

        'title': title,

        'type': ctype,

        'disability_category': disability_category,

        'text': final_text,

        'file': filename,

        'sign_video': sign_filename,

        'youtube_link': youtube_link,

        'teacher': session['name'],

        'status': status,

        'ts': time.time(),

        'transcript': final_text,

        'summary': summary,

        'braille': to_braille(final_text)

    }

    db['content'].insert(0, item)

    save_db(db)

    return redirect('/teacher')

@app.route('/teacher/delete/<cid>')
def delete(cid):
    if session.get('role') != 'teacher':
        return redirect('/login')
    db = load_db()
    item = next((c for c in db['content'] if c['id'] == cid), None)
    if item and item.get('file'):
        fpath = os.path.join(UPLOAD_FOLDER, item['file'])
        if os.path.exists(fpath):
            os.remove(fpath)
    db['content'] = [c for c in db['content'] if c['id'] != cid]
    save_db(db)
    return redirect('/teacher')

# ── Student ──────────────────────────────────────────────────────────────────

@app.route('/student')
def student():
    if session.get('role') != 'student':
        return redirect('/login')
    db = load_db()
    return render_template('student.html', content=db['content'],
                           name=session['name'], disability=session.get('disability', 'none'))

@app.route('/content/<cid>')
def view_content(cid):
    if 'uid' not in session:
        return redirect('/login')
    db = load_db()
    item = next((c for c in db['content'] if c['id'] == cid), None)
    if not item:
        return 'Not found', 404
    return render_template('content.html', item=item, disability=session.get('disability', 'none'))

@app.route('/uploads/<path:filename>')
def uploaded_file(filename):
    return send_from_directory(UPLOAD_FOLDER, filename)

# ── API: Summarize ────────────────────────────────────────────────────────────

@app.route('/api/summarize', methods=['POST'])
def summarize():
    data = request.json
    text = data.get('text', '')
    if not text.strip():
        return jsonify({'summary': 'No content to summarize.'})
    summary = summarize_as_bullets(text)
    return jsonify({'summary': summary})

# ── API: Braille ──────────────────────────────────────────────────────────────

@app.route('/api/braille', methods=['POST'])
def braille_api():
    text = request.json.get('text', '')
    return jsonify({'braille': to_braille(text)})

# ── API: Extract ──────────────────────────────────────────────────────────────

@app.route('/api/extract', methods=['POST'])
def api_extract():
    cid = request.json.get('cid', '')
    db = load_db()
    item = next((c for c in db['content'] if c['id'] == cid), None)
    if not item:
        return jsonify({'error': 'Not found'}), 404
    return jsonify({'text': item.get('text', ''), 'transcript': item.get('transcript', '')})

# ── TTS ───────────────────────────────────────────────────────────────────────

@app.route('/tts')
def tts():
    text = request.args.get('text', '')
    if not text:
        return 'No text', 400

    async def generate():
        communicate = edge_tts.Communicate(text, voice='en-US-JennyNeural')
        output = 'static/tts_output.mp3'
        await communicate.save(output)

    asyncio.run(generate())
    return send_file('static/tts_output.mp3', mimetype='audio/mpeg')

if __name__ == '__main__':
    app.run(host='0.0.0.0', debug=True, port=5000)